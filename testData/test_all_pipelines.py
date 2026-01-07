import asyncio
import os
import sys
from pathlib import Path
import json
from pprint import pprint
# Add the agent directory to sys.path so we can import my_agent
# Assuming the script is run from inside testData folder which is sibling to excel_analysis_agent
AGENT_ROOT = Path(__file__).parent.parent / "excel_analysis_agent"
sys.path.append(str(AGENT_ROOT))

try:
    from my_agent.agent import graph
    from langchain_core.messages import HumanMessage
    print("✅ Successfully imported agent graph")
except ImportError as e:
    print(f"❌ Error: Could not import my_agent. Ensure {AGENT_ROOT} is correct.")
    print(f"Details: {e}")
    sys.exit(1)

# Helper to create sample files for testing if they don't exist
def create_sample_files():
    print("\n🛠️ Checking for sample files in testData...")
    
    # Excel sample (already exists usually, but just in case)
    excel_file = "esd.xlsx"
    if not os.path.exists(excel_file):
        print(f"⚠️ Warning: {excel_file} not found. Excel tests might fail.")
    else:
        print(f"✅ Found {excel_file}")

    # Word sample
    docx_file = "sample_report.docx"
    if not os.path.exists(docx_file):
        try:
            from docx import Document
            doc = Document()
            doc.add_heading('Quarterly Sales Report', 0)
            doc.add_paragraph('This is a test document for the Multi-Asset Pipeline.')
            doc.add_paragraph('Growth was 15% in Q3, driven by the electronics sector.')
            doc.save(docx_file)
            print(f"✅ Created {docx_file}")
        except Exception as e:
            print(f"❌ Could not create {docx_file}: {e}")
    else:
        print(f"✅ Found {docx_file}")

    # Text sample
    txt_file = "notes.txt"
    if not os.path.exists(txt_file):
        with open(txt_file, "w") as f:
            f.write("Project Notes:\n- Deadline: Friday\n- Team: Analysis Group\n- Priority: High")
        print(f"✅ Created {txt_file}")
    else:
        print(f"✅ Found {txt_file}")

    # PowerPoint sample
    pptx_file = "presentation.pptx"
    if not os.path.exists(pptx_file):
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Market Trends 2024"
            subtitle.text = "Analysis performed by AI Agent"
            prs.save(pptx_file)
            print(f"✅ Created {pptx_file}")
        except Exception as e:
            print(f"❌ Could not create {pptx_file}: {e}")
    else:
        print(f"✅ Found {pptx_file}")

async def run_test(name, file_path, query):
    print(f"\n{'='*60}")
    print(f"🚀 TESTING PIPELINE: {name}")
    print(f"📂 File: {file_path}")
    print(f"🤖 Query: {query}")
    print(f"{'='*60}\n")

    if not os.path.exists(file_path):
        print(f"❌ Error: File {file_path} not found. Skipping test.")
        return

    # In LangGraph, we use the unified state
    initial_state = {
        "messages": [HumanMessage(content=query)],
        "file_path": os.path.abspath(file_path),
    }
    print("Initial State:")
    pprint(initial_state, indent=2)

    print("Running agent...")
    try:
        # Collect all events for tracking
        all_events = []
        # final_state = None
        async for event in graph.astream(initial_state, stream_mode="values"):
            # final_state = event
            all_events.append(event)  # Store the full event (state snapshot)
            
            # Existing partial print logic (optional, for real-time feedback)
            # if "messages" in event and event["messages"]:
                # last_msg = event["messages"][-1]
                # if last_msg.name == "CodingAgent":
                    # print(f"📝 Agent Response (partial): {last_msg.content[:100]}...")
        
        # After streaming, save all events to JSON
        events_filename = f"events_{name.replace(' ', '_')}.json"
        with open(events_filename, "w") as f:
            json.dump(all_events, f, indent=2, default=str)  # Use default=str to handle non-serializable objects like HumanMessage
        print(f"✅ All events saved to {events_filename}")
        
        # print(f"\n✅ {name} Test Complete!")
        # if final_state and "final_analysis" in final_state:
        #     print("\nFinal Analysis Result:")
        #     print("-" * 40)
        #     print(final_state["final_analysis"][:500] + "..." if len(final_state["final_analysis"]) > 500 else final_state["final_analysis"])
        
    except Exception as e:
        print(f"❌ Error during {name} test: {e}")

async def main():
    # Make sure we are in the correct directory
    os.chdir(Path(__file__).parent)
    
    # 1. Create dummy samples if needed
    # create_sample_files()

    # 2. Run sequential tests
    
    # Excel Test
    # await run_test(
    #     "Excel Pipeline", 
    #     "esd.xlsx", 
    #     "Give me a table of employees having salary more than $100000 and live in United States i want top 5 richest employees"
    # )

    # Document Test (Word)
    # await run_test(
    #     "Document Pipeline (Word)", 
    #     "sample_report.docx", 
    #     "What was the growth in Q3 and which sector drove it?"
    # )

    # Document Test (Text)
    await run_test(
        "Document Pipeline (Text)", 
        "notes.txt", 
        "List down all project id and also count of tasks in each project"
    )

    # PowerPoint Test
    # await run_test(
    #     "PowerPoint Pipeline", 
    #     "presentation.pptx", 
    #     "What is the title and subtitle of this presentation?"
    # )

if __name__ == "__main__":
    asyncio.run(main())
