"""PowerPoint inspection and slide extraction."""

import asyncio
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List


async def extract_slides(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract slide content from PowerPoint presentation.
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        List of slide dictionaries with content
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise ImportError(
            "python-pptx not installed. Run: pip install python-pptx"
        )
    
    def _extract():
        prs = Presentation(file_path)
        slides = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": "",
                "shapes": [],
                "has_images": False,
                "has_charts": False,
                "has_tables": False,
            }
            
            for shape in slide.shapes:
                # Check shape types
                if hasattr(shape, 'shape_type'):
                    shape_type = shape.shape_type
                    slide_data["shapes"].append(str(shape_type))
                    
                    # Detect images, charts, tables
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        slide_data["has_images"] = True
                    elif shape_type == MSO_SHAPE_TYPE.CHART:
                        slide_data["has_charts"] = True
                    elif shape_type == MSO_SHAPE_TYPE.TABLE:
                        slide_data["has_tables"] = True
                        # Extract table data
                        if hasattr(shape, 'table'):
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            slide_data["content"].append(f"[TABLE: {table_data}]")
                
                # Extract text
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        # Check if this is the title
                        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                            if hasattr(shape, 'placeholder_format'):
                                ph_type = shape.placeholder_format.type
                                # Type 1 is typically title
                                if ph_type == 1:
                                    slide_data["title"] = text
                                    continue
                        slide_data["content"].append(text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data["notes"] = notes_frame.text.strip()
            
            slides.append(slide_data)
        
        return slides
    
    return await asyncio.to_thread(_extract)


async def inspect_presentation(file_path: str) -> Dict[str, Any]:
    """
    Full inspection of a PowerPoint file.
    
    Extracts all slides, content, speaker notes, and metadata.
    Uses full-context approach - all content is passed to LLM.
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        Structured data context for the agent
    """
    slides = await extract_slides(file_path)
    file_name = Path(file_path).name
    
    # Analyze content
    total_text_parts = []
    for slide in slides:
        title = slide.get("title", "")
        content = slide.get("content", [])
        if title:
            total_text_parts.append(title)
        total_text_parts.extend(content)
    
    total_text = " ".join(total_text_parts)
    word_count = len(total_text.split())
    
    titles = [s["title"] for s in slides if s.get("title")]
    slides_with_notes = sum(1 for s in slides if s.get("notes"))
    slides_with_images = sum(1 for s in slides if s.get("has_images"))
    slides_with_charts = sum(1 for s in slides if s.get("has_charts"))
    slides_with_tables = sum(1 for s in slides if s.get("has_tables"))
    
    # Generate human-readable description
    description_parts = [
        "Presentation Overview:",
        f"- File: {file_name}",
        f"- Total Slides: {len(slides)}",
        f"- Word Count: {word_count}",
        f"- Slides with Speaker Notes: {slides_with_notes}",
    ]
    
    if slides_with_images:
        description_parts.append(f"- Slides with Images: {slides_with_images}")
    if slides_with_charts:
        description_parts.append(f"- Slides with Charts: {slides_with_charts}")
    if slides_with_tables:
        description_parts.append(f"- Slides with Tables: {slides_with_tables}")
    
    # Add slide titles
    if titles:
        description_parts.append("\nSlide Titles:")
        for i, title in enumerate(titles, 1):
            description_parts.append(f"  {i}. {title}")
    
    # Add detailed slide content
    description_parts.append("\n\nSlide Content:")
    for slide in slides:
        slide_num = slide["slide_number"]
        title = slide.get("title", "Untitled")
        content = slide.get("content", [])
        notes = slide.get("notes", "")
        
        description_parts.append(f"\n--- Slide {slide_num}: {title} ---")
        if content:
            for item in content:
                description_parts.append(f"  • {item[:200]}{'...' if len(item) > 200 else ''}")
        if notes:
            description_parts.append(f"  [Speaker Notes: {notes[:150]}{'...' if len(notes) > 150 else ''}]")
    
    description = "\n".join(description_parts)
    
    # Create full text for LLM context
    full_text_parts = []
    for slide in slides:
        slide_num = slide["slide_number"]
        title = slide.get("title", "Untitled")
        content = slide.get("content", [])
        notes = slide.get("notes", "")
        
        full_text_parts.append(f"=== SLIDE {slide_num}: {title} ===")
        full_text_parts.extend(content)
        if notes:
            full_text_parts.append(f"[SPEAKER NOTES: {notes}]")
        full_text_parts.append("")
    
    full_text = "\n".join(full_text_parts)
    
    return {
        "file_path": os.path.abspath(file_path),
        "file_name": file_name,
        "analyzed_at": datetime.now().isoformat(),
        "description": description,
        "full_text": full_text,  # Full content for LLM context
        "slides": slides,
        "summary": {
            "slide_count": len(slides),
            "word_count": word_count,
            "has_notes": slides_with_notes > 0,
            "has_images": slides_with_images > 0,
            "has_charts": slides_with_charts > 0,
            "has_tables": slides_with_tables > 0,
            "titles": titles,
        },
        "capabilities": [
            "slide_search",
            "summarization",
            "content_extraction",
            "notes_extraction",
            "presentation_analysis",
        ]
    }
